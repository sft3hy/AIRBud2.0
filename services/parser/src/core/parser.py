import os
import io
import fitz  # PyMuPDF
import cv2
from PIL import Image
from docx import Document
from pptx import Presentation
from typing import List, Tuple, Dict, Any, Generator
import subprocess
import pandas as pd
import tempfile
import shutil
import time
from moviepy import VideoFileClip

from src.utils.logger import logger
from src.core.detector import ChartDetector


class DocumentParser:
    def __init__(self, detector: ChartDetector, output_dir: str):
        self.output_dir = output_dir
        self.detector = detector

    def parse(self, file_path: str) -> Generator[Dict[str, Any], None, None]:
        """
        Yields progress updates and finally the result.
        Yield schema:
        - {"status": "processing", "step": "...", "progress": float}
        - {"status": "completed", "result": {...}}
        - {"status": "error", "error": "..."}
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        file_ext = os.path.splitext(file_path)[1].lower()
        file_name = os.path.splitext(os.path.basename(file_path))[0]

        # Create a specific folder for this file's assets
        self.file_output_dir = os.path.join(self.output_dir, file_name)
        os.makedirs(self.file_output_dir, exist_ok=True)

        logger.info(f"Parsing {file_ext} document: {file_name}")
        start_time = time.time()

        text = ""
        images = []
        audio_path = None

        try:
            yield {
                "status": "processing", 
                "step": f"Initializing {file_ext} parser", 
                "progress": 0.0,
                "start_time": start_time
            }

            if file_ext == ".pdf":
                for update in self._extract_from_pdf(file_path):
                    if isinstance(update, dict):
                        # Add timing info to every update
                        update["elapsed"] = time.time() - start_time
                        yield update
                    else:
                        # Final result from sub-generator
                        text, images = update
            
            elif file_ext == ".docx":
                text, images = self._extract_from_docx(file_path)
            
            elif file_ext == ".pptx":
                for update in self._extract_from_pptx(file_path):
                    if isinstance(update, dict):
                        update["elapsed"] = time.time() - start_time
                        yield update
                    else:
                        text, images = update
            
            elif file_ext == ".xlsx":
                text = self._extract_from_xlsx(file_path)
            
            elif file_ext == ".txt":
                text = self._extract_from_txt(file_path)
            
            elif file_ext == ".mp4":
                # Video parsing might take a while, let's wrap it if needed or just yield one start event
                # For now, we'll keep it simple as it's not the primary focus of the change, 
                # but adding a yield here would be good.
                yield {"status": "processing", "step": "Processing video...", "progress": 0.1}
                text, images, audio_path = self._extract_from_mp4(file_path, file_name)
            
            else:
                raise ValueError(f"Unsupported format: {file_ext}")
            
            end_time = time.time()
            duration = end_time - start_time
            
            yield {
                "status": "completed",
                "result": {
                    "text": text,
                    "images": images,
                    "audio_path": audio_path
                },
                "metrics": {
                    "duration": duration,
                    "images_found": len(images)
                }
            }

        except Exception as e:
            logger.error(f"Error parsing {file_path}: {e}", exc_info=True)
            yield {"status": "error", "error": str(e)}
            # We don't raise here to allow the stream to close gracefully with the error message

    def _extract_from_xlsx(self, path: str) -> str:
        full_text = []
        try:
            xls = pd.read_excel(path, sheet_name=None)
            for sheet_name, df in xls.items():
                full_text.append(f"## Sheet: {sheet_name}\n")
                df = df.fillna("")
                markdown_table = df.to_markdown(index=False)
                full_text.append(markdown_table)
                full_text.append("\n\n")
        except Exception as e:
            logger.error(f"Excel extraction failed: {e}")
            raise ValueError(f"Could not parse Excel file: {e}")
        return "".join(full_text)

    def _extract_from_pdf(self, path: str) -> Generator[Any, None, None]:
        doc = fitz.open(path)
        full_text = []
        extracted_images = []
        total_pages = len(doc)

        yield {"status": "processing", "step": f"PDF loaded ({total_pages} pages)", "progress": 0.05}

        for i, page in enumerate(doc):
            # 1. Extract Text
            full_text.append(f"## Page {i+1}\n{page.get_text()}")

            # 2. Render Page for Vision Detection
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            img = Image.open(io.BytesIO(pix.tobytes("png")))

            # 3. Detect & Crop
            crops = self._process_visuals(img, f"page{i+1}")
            extracted_images.extend(crops)

            # 4. Insert Placeholders
            for crop_path in crops:
                filename = os.path.basename(crop_path)
                full_text.append(f"\n[CHART_PLACEHOLDER:{filename}]\n")
            
            # Yield Progress
            progress = (i + 1) / total_pages
            yield {
                "status": "processing", 
                "step": f"Processed Page {i+1}/{total_pages}", 
                "progress": progress,
                "current_images": len(extracted_images)
            }

        doc.close()
        yield "\n".join(full_text), extracted_images

    def _extract_from_docx(self, path: str) -> Tuple[str, List[str]]:
        # Non-streaming implementation for now (can be upgraded later)
        doc = Document(path)
        full_text = []
        extracted_images = []

        for para in doc.paragraphs:
            full_text.append(para.text)

        count = 0
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                try:
                    img_data = rel.target_part.blob
                    img = Image.open(io.BytesIO(img_data))
                    if img.width > 200 and img.height > 200:
                        count += 1
                        fname = f"docx_visual_{count}.png"
                        save_path = os.path.join(self.file_output_dir, fname)
                        img.save(save_path)
                        extracted_images.append(save_path)
                        full_text.append(f"\n[CHART_PLACEHOLDER:{fname}]\n")
                except Exception as e:
                    logger.warning(f"Failed to extract DOCX image: {e}")

        return "\n\n".join(full_text), extracted_images

    def _extract_from_pptx(self, path: str) -> Generator[Any, None, None]:
        prs = Presentation(path)
        full_text = []
        extracted_images = []
        
        # 1. Convert Slides to Images
        yield {"status": "processing", "step": "Converting slides...", "progress": 0.1}
        slide_images = self._convert_pptx_to_images(path)
        total_slides = len(prs.slides)
        
        yield {"status": "processing", "step": f"Slides converted ({len(slide_images)}). Extracting content...", "progress": 0.2}

        for i, slide in enumerate(prs.slides):
            full_text.append(f"## Slide {i+1}")

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    full_text.append(shape.text)

            if i < len(slide_images):
                crops = self._process_visuals(slide_images[i], f"slide{i+1}")
                extracted_images.extend(crops)

                for crop_path in crops:
                    filename = os.path.basename(crop_path)
                    full_text.append(f"\n[CHART_PLACEHOLDER:{filename}]\n")
            
            progress = 0.2 + (0.8 * ((i + 1) / total_slides))
            yield {
                "status": "processing",
                "step": f"Processed Slide {i+1}/{total_slides}",
                "progress": progress,
                "current_images": len(extracted_images)
            }

        yield "\n\n".join(full_text), extracted_images

    def _convert_pptx_to_images(self, pptx_path: str) -> List[Image.Image]:
        images = []
        try:
            with tempfile.TemporaryDirectory() as tmpdir:
                cmd = ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmpdir, pptx_path]
                subprocess.run(cmd, check=True, timeout=120, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
                
                pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
                pdf_path = os.path.join(tmpdir, pdf_name)

                if os.path.exists(pdf_path):
                    doc = fitz.open(pdf_path)
                    for page in doc:
                        pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                        img = Image.open(io.BytesIO(pix.tobytes("png")))
                        images.append(img)
                    doc.close()
        except Exception as e:
            logger.error(f"PPTX conversion failed: {e}")
        return images
    
    def _extract_from_txt(self, path: str) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_from_mp4(self, path: str, filename_base: str) -> Tuple[str, List[str], str]:
        images = []
        full_text = [f"# Video Analysis: {filename_base}\n"]
        audio_path = os.path.join(self.file_output_dir, f"{filename_base}.mp3")
        
        try:
            clip = VideoFileClip(path)
            if clip.audio:
                clip.audio.write_audiofile(audio_path, logger=None)
            else:
                audio_path = None
            clip.close()
        except Exception as e:
            logger.error(f"Audio extraction failed: {e}")
            audio_path = None

        cap = cv2.VideoCapture(path)
        fps = cap.get(cv2.CAP_PROP_FPS) or 24
        frame_interval = int(fps * 5)
        max_frames = 20
        count = 0
        saved_count = 0

        while cap.isOpened() and saved_count < max_frames:
            ret, frame = cap.read()
            if not ret: break
            
            if count % frame_interval == 0:
                rgb_frame = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                img = Image.fromarray(rgb_frame)
                timestamp = count / fps
                fname = f"frame_{int(timestamp)}s.png"
                save_path = os.path.join(self.file_output_dir, fname)
                img.save(save_path)
                images.append(save_path)
                full_text.append(f"\n## Timestamp: {int(timestamp)}s\n[CHART_PLACEHOLDER:{fname}]\n")
                saved_count += 1
            count += 1
        cap.release()
        
        return "\n".join(full_text), images, audio_path

    def _process_visuals(self, page_image: Image.Image, prefix: str) -> List[str]:
        bboxes = self.detector.detect(page_image)
        saved_paths = []
        for i, (x1, y1, x2, y2) in enumerate(bboxes):
            try:
                crop = page_image.crop((x1, y1, x2, y2))
                fname = f"{prefix}_visual_{i+1}.png"
                path = os.path.join(self.file_output_dir, fname)
                crop.save(path)
                saved_paths.append(path)
            except Exception as e:
                logger.warning(f"Failed to save crop {i}: {e}")
        return saved_paths

