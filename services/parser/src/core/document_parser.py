import os
import re
import io
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from pptx import Presentation
from typing import List, Tuple, Dict, Optional
import subprocess

# Import Chart Detector
from src.utils.chart_detection import PubLayNetDetector


class DocumentParser:
    def __init__(self, vision_model, output_dir: str):
        self.output_dir = output_dir
        self.layout_detector = PubLayNetDetector(confidence_threshold=0.5, padding=60)

    def parse_and_get_images(self, file_path: str) -> Tuple[str, List[str]]:
        """
        Parses document, extracts text, detects charts, saves crops.
        Returns: (markdown_text, list_of_saved_image_paths)
        """
        file_ext = os.path.splitext(file_path)[1].lower()
        file_output_dir = os.path.join(
            self.output_dir, os.path.splitext(os.path.basename(file_path))[0]
        )
        os.makedirs(file_output_dir, exist_ok=True)

        extracted_images = []

        if file_ext == ".pdf":
            text = self._extract_from_pdf(file_path, file_output_dir, extracted_images)
        elif file_ext == ".docx":
            text = self._extract_from_docx(file_path, file_output_dir, extracted_images)
        elif file_ext == ".pptx":
            text = self._extract_from_pptx(file_path, file_output_dir, extracted_images)
        else:
            raise ValueError(f"Unsupported format: {file_ext}")

        return text, extracted_images

    def _extract_from_pdf(self, path, output_dir, img_list):
        doc = fitz.open(path)
        full_text = []
        for i, page in enumerate(doc):
            # Text
            full_text.append(f"## Page {i+1}\n{page.get_text()}")

            # Image for detection
            pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
            img = Image.open(io.BytesIO(pix.tobytes("png")))

            # Detect & Crop
            crops = self._process_visuals(img, f"page{i+1}", output_dir)
            img_list.extend(crops)

            # Add placeholders
            for crop_path in crops:
                filename = os.path.basename(crop_path)
                full_text.append(f"\n[CHART_PLACEHOLDER:{filename}]\n")

        return "\n".join(full_text)

    def _extract_from_docx(self, path, output_dir, img_list):
        doc = Document(path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)

        # Extract images from relationships
        for rel in doc.part.rels.values():
            if "image" in rel.target_ref:
                img_data = rel.target_part.blob
                img = Image.open(io.BytesIO(img_data))
                if img.width > 150 and img.height > 150:
                    fname = f"docx_img_{len(img_list)}.png"
                    save_path = os.path.join(output_dir, fname)
                    img.save(save_path)
                    img_list.append(save_path)
                    full_text.append(f"\n[CHART_PLACEHOLDER:{fname}]\n")
        return "\n".join(full_text)

    def _extract_from_pptx(self, path, output_dir, img_list):
        # Convert to PDF first for layout detection (simplified approach)
        # Or iterate slides
        prs = Presentation(path)
        full_text = []

        # Note: Proper PPTX chart extraction requires rendering slides to images
        # We will use the libreoffice trick from original code if available,
        # otherwise just text. For stability in docker, we might skip the libreoffice dependency
        # unless strictly needed, but let's include text extraction.

        for i, slide in enumerate(prs.slides):
            full_text.append(f"## Slide {i+1}")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)

        return "\n".join(full_text)

    def _process_visuals(self, page_image, prefix, output_dir) -> List[str]:
        bboxes = self.layout_detector.detect(page_image)
        saved_paths = []
        for i, (x1, y1, x2, y2) in enumerate(bboxes):
            crop = page_image.crop((x1, y1, x2, y2))
            fname = f"{prefix}_visual_{i+1}.png"
            path = os.path.join(output_dir, fname)
            crop.save(path)
            saved_paths.append(path)
        return saved_paths
